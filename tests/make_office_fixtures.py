"""Generate the golden Office/ODF round-trip fixtures into tests/data/input/.

These are tiny, self-authored, single-page documents — copyright is ours (shipped
under the repo licence). Each embeds a whitespace-free sentinel (see
``office_fixtures.SENTINELS``) and has **no external references** (attached
template / linked content / external data), so the OOXML ``.rels`` pre-scan never
refuses them and ``office_via_com`` converts them deterministically to a 1-page PDF.

One-time / regenerate only — the produced files are committed, so ordinary
contributors never need the generator libs. To (re)generate:

    pip install python-docx openpyxl python-pptx odfpy
    python tests/make_office_fixtures.py
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from office_fixtures import INPUT_DIR, SENTINELS  # noqa: E402


def make_docx(path: Path, sentinel: str) -> None:
    from docx import Document
    d = Document()  # default Normal template — no attached template
    d.add_paragraph(sentinel)
    d.save(str(path))


def make_xlsx(path: Path, sentinel: str) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws["A1"] = sentinel
    # Widen A so the sentinel isn't clipped to the column on PDF export (a narrow
    # column truncates the rendered text even though the cell value is intact).
    ws.column_dimensions["A"].width = len(sentinel) + 5
    ws.print_area = "A1"  # one cell -> one printed page
    wb.save(str(path))


def make_pptx(path: Path, sentinel: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = Blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text_frame.text = sentinel
    prs.save(str(path))


def make_odt(path: Path, sentinel: str) -> None:
    from odf.opendocument import OpenDocumentText
    from odf.text import P
    doc = OpenDocumentText()
    doc.text.addElement(P(text=sentinel))
    doc.save(str(path))


def make_ods(path: Path, sentinel: str) -> None:
    from odf.opendocument import OpenDocumentSpreadsheet
    from odf.table import Table, TableRow, TableCell
    from odf.text import P
    doc = OpenDocumentSpreadsheet()
    table = Table(name="Sheet1")
    row = TableRow()
    cell = TableCell(valuetype="string")
    cell.addElement(P(text=sentinel))
    row.addElement(cell)
    table.addElement(row)
    doc.spreadsheet.addElement(table)
    doc.save(str(path))


def make_odp(path: Path, sentinel: str) -> None:
    from odf.opendocument import OpenDocumentPresentation
    from odf.style import MasterPage, PageLayout, PageLayoutProperties
    from odf.draw import Page, Frame, TextBox
    from odf.text import P
    doc = OpenDocumentPresentation()
    layout = PageLayout(name="PL")
    doc.automaticstyles.addElement(layout)
    layout.addElement(PageLayoutProperties(
        pagewidth="25.4cm", pageheight="19.05cm", printorientation="landscape"))
    master = MasterPage(name="Master", pagelayoutname=layout)
    doc.masterstyles.addElement(master)
    page = Page(masterpagename=master)
    doc.presentation.addElement(page)
    frame = Frame(width="20cm", height="5cm", x="2cm", y="7cm")
    page.addElement(frame)
    box = TextBox()
    frame.addElement(box)
    box.addElement(P(text=sentinel))
    doc.save(str(path))


_MAKERS = {
    ".docx": make_docx, ".xlsx": make_xlsx, ".pptx": make_pptx,
    ".odt": make_odt, ".ods": make_ods, ".odp": make_odp,
}


def main() -> None:
    INPUT_DIR.mkdir(parents=True, exist_ok=True)
    for fname, sentinel in SENTINELS.items():
        path = INPUT_DIR / fname
        _MAKERS[path.suffix.lower()](path, sentinel)
        print(f"wrote {path}  ({path.stat().st_size} B)  sentinel={sentinel}")


if __name__ == "__main__":
    main()
